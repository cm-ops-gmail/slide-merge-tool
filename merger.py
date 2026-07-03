"""
merger.py — Core PPTX merge logic.

Takes a root PPTX (content) and a template PPTX (design) as bytes,
returns a new styled PPTX as bytes.
"""

import io
import copy
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


# ─────────────────────────────────────────────
# Theme-colour resolution
# ─────────────────────────────────────────────

_THEME_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


def _build_theme_map(prs):
    """Map every MSO_THEME_COLOR to a concrete RGBColor from *this* file's theme.

    Root and Template can carry different colour schemes, so a scheme-coloured
    run/shape must be resolved to RGB against the source (root) theme and baked
    in — otherwise it would re-resolve against the destination theme and change
    colour. Returns {} if the theme can't be read (callers then no-op).
    """
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(_THEME_RELTYPE)
        root = etree.fromstring(theme_part.blob)
        clr_scheme = root.find(qn('a:themeElements')).find(qn('a:clrScheme'))
        slot = {}
        for child in clr_scheme:
            tag = child.tag.split('}')[1]          # dk1, lt1, accent1, …
            srgb = child.find(qn('a:srgbClr'))
            if srgb is not None:
                slot[tag] = srgb.get('val')
            else:
                sysclr = child.find(qn('a:sysClr'))
                slot[tag] = sysclr.get('lastClr') if sysclr is not None else '000000'

        clr_map_el = master.element.find(qn('p:clrMap'))
        clr_map = dict(clr_map_el.attrib) if clr_map_el is not None else {}

        def hexof(name):
            return slot.get(name, '000000')

        M = MSO_THEME_COLOR
        mapping = {
            M.DARK_1:  hexof('dk1'),      M.LIGHT_1: hexof('lt1'),
            M.DARK_2:  hexof('dk2'),      M.LIGHT_2: hexof('lt2'),
            M.ACCENT_1: hexof('accent1'), M.ACCENT_2: hexof('accent2'),
            M.ACCENT_3: hexof('accent3'), M.ACCENT_4: hexof('accent4'),
            M.ACCENT_5: hexof('accent5'), M.ACCENT_6: hexof('accent6'),
            M.HYPERLINK: hexof('hlink'),  M.FOLLOWED_HYPERLINK: hexof('folHlink'),
            M.TEXT_1: hexof(clr_map.get('tx1', 'dk1')),
            M.BACKGROUND_1: hexof(clr_map.get('bg1', 'lt1')),
            M.TEXT_2: hexof(clr_map.get('tx2', 'dk2')),
            M.BACKGROUND_2: hexof(clr_map.get('bg2', 'lt2')),
        }
        return {k: RGBColor.from_string(v) for k, v in mapping.items()}
    except Exception:
        return {}


def _apply_brightness(rgb, brightness):
    """Approximate a scheme colour's lumMod/lumOff tint on a concrete RGB."""
    if not brightness:
        return rgb
    r, g, b = rgb[0], rgb[1], rgb[2]
    if brightness > 0:      # tint toward white
        r = int(round(r + (255 - r) * brightness))
        g = int(round(g + (255 - g) * brightness))
        b = int(round(b + (255 - b) * brightness))
    else:                   # shade toward black
        f = 1 + brightness
        r, g, b = int(round(r * f)), int(round(g * f)), int(round(b * f))
    clamp = lambda v: max(0, min(255, v))
    return RGBColor(clamp(r), clamp(g), clamp(b))


def _resolve_color(color, theme_map):
    """Return an RGBColor for a run/fill ColorFormat, or None to leave as-is."""
    try:
        if color is None:
            return None
        if color.type == 1:                         # already RGB
            return color.rgb
        if color.type == 2 and theme_map:           # scheme/theme colour
            base = theme_map.get(color.theme_color)
            if base is None:
                return None
            try:
                return _apply_brightness(base, color.brightness)
            except Exception:
                return base
    except Exception:
        return None
    return None


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def _is_background_shape(shape):
    """Skip large full-width background rectangles."""
    if shape.shape_type == 1:  # AUTO_SHAPE
        if shape.left <= 100_000 and shape.width >= 9_000_000:
            return True
    return False


def _copy_text_frame(src_tf, dst_tf, font_name="Noto Sans Bengali", theme_map=None):
    """Copy every paragraph & run.

    If ``font_name`` is given, every run is normalised to that font;
    if it is ``None`` the run's original font name is preserved.
    ``theme_map`` resolves scheme/theme font colours to concrete RGB.
    """
    dst_tf.clear()
    dst_tf.word_wrap = src_tf.word_wrap

    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        try:
            setattr(dst_tf, attr, getattr(src_tf, attr))
        except Exception:
            pass

    for p_idx, src_p in enumerate(src_tf.paragraphs):
        dst_p = dst_tf.paragraphs[0] if p_idx == 0 else dst_tf.add_paragraph()

        try:
            dst_p.alignment = src_p.alignment
        except Exception:
            pass
        for attr in ("space_after", "space_before", "line_spacing", "level"):
            try:
                setattr(dst_p, attr, getattr(src_p, attr))
            except Exception:
                pass

        for src_run in src_p.runs:
            dst_run = dst_p.add_run()
            dst_run.text = src_run.text
            f = dst_run.font
            if font_name is not None:
                f.name = font_name
            elif src_run.font.name:
                f.name = src_run.font.name
            if src_run.font.size:
                f.size = src_run.font.size
            f.bold      = src_run.font.bold
            f.italic    = src_run.font.italic
            f.underline = src_run.font.underline
            rgb = _resolve_color(src_run.font.color, theme_map)
            if rgb is not None:
                try:
                    f.color.rgb = rgb
                except Exception:
                    pass


def _add_reading_tag(slide):
    """Grey badge + red 'Reading' label in the top-left of every content slide."""
    try:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 457_200, 514_350, 1_554_480, 347_472)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xD8, 0xD8, 0xD8)
        box.line.fill.background()
    except Exception:
        pass

    try:
        tb = slide.shapes.add_textbox(533_400, 529_469, 1_449_177, 347_472)
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = "Reading"
        run.font.name = "Noto Sans Bengali"
        run.font.size = 228_600
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xC3, 0x23, 0x26)
    except Exception:
        pass


def _add_title_header(slide, title_tf):
    """Styled slide title + red divider line for standard content slides."""
    try:
        title_box = slide.shapes.add_textbox(457_200, 960_120, 8_229_600, 411_480)
        _copy_text_frame(title_tf, title_box.text_frame)
        p = title_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = "Noto Sans Bengali"
            run.font.size = 304_800
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    except Exception:
        pass

    try:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 457_200, 1_385_316, 8_229_600, 22_860)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
        line.line.fill.background()
    except Exception:
        pass


def _copy_table(new_slide, src_shape, left, top, width, height,
                font_name="Noto Sans Bengali", theme_map=None):
    rows = len(src_shape.table.rows)
    cols = len(src_shape.table.columns)
    new_tbl_shape = new_slide.shapes.add_table(rows, cols, left, top, width, height)
    new_tbl = new_tbl_shape.table

    # Column / row sizes
    for ci, col in enumerate(src_shape.table.columns):
        new_tbl.columns[ci].width = col.width
    for ri, row in enumerate(src_shape.table.rows):
        new_tbl.rows[ri].height = row.height

    # Merges first
    for ri in range(rows):
        for ci in range(cols):
            src_cell = src_shape.table.cell(ri, ci)
            if src_cell.is_merge_origin and (src_cell.span_height > 1 or src_cell.span_width > 1):
                tr = ri + src_cell.span_height - 1
                tc = ci + src_cell.span_width - 1
                new_tbl.cell(ri, ci).merge(new_tbl.cell(tr, tc))

    # Content
    for ri in range(rows):
        for ci in range(cols):
            src_cell = src_shape.table.cell(ri, ci)
            if not src_cell.is_spanned or src_cell.is_merge_origin:
                dst_cell = new_tbl.cell(ri, ci)
                _copy_text_frame(src_cell.text_frame, dst_cell.text_frame, font_name, theme_map)
                try:
                    if src_cell.fill.type == 1:
                        rgb = _resolve_color(src_cell.fill.fore_color, theme_map)
                        if rgb is not None:
                            dst_cell.fill.solid()
                            dst_cell.fill.fore_color.rgb = rgb
                except Exception:
                    pass


def _copy_shape(new_slide, shape, left, top, width, height,
                font_name="Noto Sans Bengali", theme_map=None):
    """Recreate a non-table, non-image shape."""
    if shape.shape_type in (17, 14):   # TEXT_BOX / PLACEHOLDER
        new_shape = new_slide.shapes.add_textbox(left, top, width, height)
    elif shape.shape_type == 1:        # AUTO_SHAPE — keep its real geometry
        try:
            new_shape = new_slide.shapes.add_shape(
                shape.auto_shape_type, left, top, width, height
            )
            # Preserve adjustments (e.g. rounded-rectangle corner radius).
            try:
                for i, adj in enumerate(shape.adjustments):
                    new_shape.adjustments[i] = adj
            except Exception:
                pass
        except Exception:
            new_shape = new_slide.shapes.add_textbox(left, top, width, height)
    else:
        try:
            new_shape = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
        except Exception:
            new_shape = new_slide.shapes.add_textbox(left, top, width, height)

    if shape.has_text_frame:
        _copy_text_frame(shape.text_frame, new_shape.text_frame, font_name, theme_map)

    try:
        if shape.fill.type == 1:
            rgb = _resolve_color(shape.fill.fore_color, theme_map)
            if rgb is not None:
                new_shape.fill.solid()
                new_shape.fill.fore_color.rgb = rgb
        elif shape.fill.type == 5:
            new_shape.fill.background()
    except Exception:
        pass

    try:
        rgb = _resolve_color(shape.line.color, theme_map)
        if rgb is not None:
            new_shape.line.color.rgb = rgb
        if shape.line.width:
            new_shape.line.width = shape.line.width
    except Exception:
        pass

    try:
        new_shape.rotation = shape.rotation
    except Exception:
        pass

    return new_shape


# ─────────────────────────────────────────────
# Main entry point
# ─────────────────────────────────────────────

def merge_presentations(root_bytes: bytes, template_bytes: bytes) -> bytes:
    """
    Merge root content into template design.

    Parameters
    ----------
    root_bytes     : bytes of the Root PPTX file
    template_bytes : bytes of the Template PPTX file

    Returns
    -------
    bytes of the new styled PPTX
    """
    prs_root     = Presentation(io.BytesIO(root_bytes))
    prs_styled   = Presentation(io.BytesIO(template_bytes))

    # ── Wipe the styled copy, but KEEP its fixed first slide ────────────────
    # The template's first slide is a fixed, designed cover (logo + branding
    # drawn directly on the slide, not reproducible from a layout), so it must
    # be preserved as the leading cover of the output rather than wiped.
    sldIdLst = prs_styled.slides._sldIdLst
    keep_first = len(sldIdLst) > 0
    start = 1 if keep_first else 0
    while len(sldIdLst) > start:
        prs_styled.part.drop_rel(sldIdLst[start].rId)
        del sldIdLst[start]

    layouts = prs_styled.slide_layouts
    # Layout index: 0=BLANK 1=TITLE 2=2_Blank 3=1_Blank 4=3_Blank 5=4_Blank 6=1_IELTS

    # Root's own theme, so scheme-coloured text/fills resolve to the colours the
    # author intended (the template may use a different scheme).
    theme_map = _build_theme_map(prs_root)

    warnings = []

    for slide_idx, root_slide in enumerate(prs_root.slides):

        # ── 1. Pick layout ────────────────────────────────────────────────────
        if slide_idx in (0, 5, 32):
            layout = layouts[1]          # TITLE
        elif (6 <= slide_idx <= 29) or (33 <= slide_idx <= 36):
            layout = layouts[4]          # 3_Blank  (passage / question slides)
        else:
            layout = layouts[2]          # 2_Blank  (standard content)

        new_slide = prs_styled.slides.add_slide(layout)

        # ── 2. Cover slide ────────────────────────────────────────────────────
        if slide_idx == 0:
            tb = new_slide.shapes.add_textbox(1_828_800, 952_189, 5_486_400, 1_619_561)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True

            data = [
                ("Class - 09",          508_000, RGBColor(0x33,0x33,0x33)),
                ("Tables and Diagrams", 508_000, RGBColor(0x33,0x33,0x33)),
                ("Reading",             304_800, RGBColor(0xC3,0x23,0x26)),
            ]
            for i, (text, size, color) in enumerate(data):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = text
                run.font.name  = "Noto Sans Bengali"
                run.font.size  = size
                run.font.bold  = True
                run.font.color.rgb = color
            continue

        # ── 3. Section-break slides ───────────────────────────────────────────
        if slide_idx in (5, 32):
            txt = "Guided Practice Questions"
            for s in root_slide.shapes:
                if s.has_text_frame and s.text_frame.text.strip() and s.left > 100_000:
                    txt = s.text_frame.text.strip()
                    break

            tb = new_slide.shapes.add_textbox(1_828_800, 1_500_000, 5_486_400, 1_619_561)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = txt
            run.font.name  = "Noto Sans Bengali"
            run.font.size  = 406_400
            run.font.bold  = True
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            continue

        # ── 4. All other slides ───────────────────────────────────────────────
        _add_reading_tag(new_slide)

        is_passage = (6 <= slide_idx <= 29) or (33 <= slide_idx <= 36)

        # Find title shape (topmost short text on non-passage slides)
        title_shape = None
        if not is_passage:
            candidates = [
                s for s in root_slide.shapes
                if s.has_text_frame
                and s.text_frame.text.strip()
                and not _is_background_shape(s)
            ]
            candidates.sort(key=lambda s: s.top)
            if candidates:
                c = candidates[0]
                if c.top < 1_000_000 and len(c.text_frame.text.strip()) < 50:
                    title_shape = c

        if title_shape:
            _add_title_header(new_slide, title_shape.text_frame)

        # Copy all remaining shapes
        for shape in root_slide.shapes:
            if _is_background_shape(shape):
                continue
            if title_shape and shape == title_shape:
                continue

            left, top, width, height = shape.left, shape.top, shape.width, shape.height

            # Shift overlapping elements down
            if is_passage:
                if left < 3_000_000 and top < 960_120:
                    shift = 960_120 - top
                    top    = 960_120
                    height = max(100_000, height - shift)
            else:
                if top < 1_500_000:
                    shift = 1_500_000 - top
                    top    = 1_500_000
                    height = max(100_000, height - shift)

            try:
                if shape.shape_type == 13:          # PICTURE
                    stream = io.BytesIO(shape.image.blob)
                    img = new_slide.shapes.add_picture(stream, left, top, width, height)
                    try:
                        img.rotation = shape.rotation
                    except Exception:
                        pass

                elif shape.shape_type == 19:        # TABLE
                    _copy_table(new_slide, shape, left, top, width, height,
                                theme_map=theme_map)

                else:
                    _copy_shape(new_slide, shape, left, top, width, height,
                                theme_map=theme_map)

            except Exception as e:
                warnings.append(f"Slide {slide_idx+1} | shape '{shape.name}': {e}")

    out = io.BytesIO()
    prs_styled.save(out)
    return out.getvalue(), warnings


# ─────────────────────────────────────────────
# SEJ merge (Spoken English Junior)
# ─────────────────────────────────────────────

def _next_shape_id(new_slide):
    """Highest cNvPr id already on the slide, + 1 (keeps shape ids unique)."""
    max_id = 1
    for el in new_slide.shapes._spTree.iter():
        if el.tag.endswith('}cNvPr'):
            try:
                max_id = max(max_id, int(el.get('id')))
            except (TypeError, ValueError):
                pass
    return max_id + 1


def _clone_shape_xml(new_slide, shape, new_id):
    """Deep-copy a shape's raw XML onto ``new_slide``, giving it a fresh id.

    Preserves geometry, theme/scheme colours, fills, outlines, connectors and
    text formatting exactly — everything the shape-rebuilding path would flatten.
    Only safe when source and destination share the same theme (they do for SEJ).
    """
    el = copy.deepcopy(shape._element)
    for cNvPr in el.iter():
        if cNvPr.tag.endswith('}cNvPr'):
            cNvPr.set('id', str(new_id))
            break
    new_slide.shapes._spTree.append(el)


def merge_sej(root_bytes: bytes, template_bytes: bytes) -> bytes:
    """
    Merge an SEJ Root file's content onto the SEJ Intermediate Template design.

    SEJ content is English, so the IELTS slide-index heuristics and Bengali
    font-normalisation are dropped. Because the Root and Template share an
    identical colour theme, every root shape is cloned at the XML level so its
    original geometry (rounded rectangles, connector "matching" lines),
    theme/scheme colours, fills and fonts survive intact. Each root slide is
    rebuilt on top of the template's design layouts (cover vs. content), so the
    template background/design shows through behind the copied content.

    Parameters
    ----------
    root_bytes     : bytes of the SEJ Root PPTX file
    template_bytes : bytes of the SEJ Intermediate Template PPTX file

    Returns
    -------
    (bytes of the new styled PPTX, list of warnings)
    """
    prs_root   = Presentation(io.BytesIO(root_bytes))
    prs_styled = Presentation(io.BytesIO(template_bytes))

    # ── Wipe all slides in the styled copy, keep the design layouts ─────────
    sldIdLst = prs_styled.slides._sldIdLst
    while len(sldIdLst):
        prs_styled.part.drop_rel(sldIdLst[0].rId)
        del sldIdLst[0]

    layouts = prs_styled.slide_layouts
    # Layout 1 = TITLE_2 (cover design), Layout 2 = TITLE_2_1 (content design).
    title_layout   = layouts[1] if len(layouts) > 1 else layouts[0]
    content_layout = layouts[2] if len(layouts) > 2 else title_layout

    warnings = []

    for slide_idx, root_slide in enumerate(prs_root.slides):
        layout = title_layout if slide_idx == 0 else content_layout
        new_slide = prs_styled.slides.add_slide(layout)

        # Copy shapes in document order so z-ordering is preserved.
        for shape in root_slide.shapes:
            try:
                if shape.shape_type == 13:          # PICTURE — re-embed the blob
                    stream = io.BytesIO(shape.image.blob)
                    img = new_slide.shapes.add_picture(
                        stream, shape.left, shape.top, shape.width, shape.height
                    )
                    try:
                        img.rotation = shape.rotation
                    except Exception:
                        pass
                else:                               # everything else — clone XML
                    # Recompute the free id each time so it can't collide with
                    # ids that add_picture() assigns itself.
                    _clone_shape_xml(new_slide, shape, _next_shape_id(new_slide))

            except Exception as e:
                warnings.append(f"Slide {slide_idx+1} | shape '{shape.name}': {e}")

    out = io.BytesIO()
    prs_styled.save(out)
    return out.getvalue(), warnings
