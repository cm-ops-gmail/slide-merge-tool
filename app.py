"""
Slide Merge Tool — Streamlit Web App
Upload a Root PPTX and a Template PPTX, merge them, download the result.
"""

import streamlit as st
from merger import merge_presentations, merge_sej

# Merge function used for each program option.
MERGE_FUNCS = {
    "IELTS": merge_presentations,
    "SEJ":   merge_sej,
}

# ─── Page config ────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Slide Merge Tool",
    page_icon="📊",
    layout="centered",
)

# ─── Custom CSS ──────────────────────────────────────────────────────────────
st.markdown("""
<style>
  /* Google Fonts */
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

  html, body, [class*="css"] { font-family: 'Inter', sans-serif; }

  /* App background */
  .stApp { background: #0f1117; color: #f1f5f9; }

  /* Hide default Streamlit chrome */
  header[data-testid="stHeader"] { display: none; }
  footer { display: none; }

  /* Hero */
  .hero {
    background: linear-gradient(135deg, #1e293b 0%, #0f172a 100%);
    border: 1px solid #334155;
    border-radius: 16px;
    padding: 40px 32px 32px;
    margin-bottom: 32px;
    text-align: center;
  }
  .hero h1 {
    font-size: 2rem; font-weight: 700;
    background: linear-gradient(90deg, #60a5fa, #a78bfa);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    margin-bottom: 8px;
  }
  .hero p { color: #94a3b8; font-size: 1rem; margin: 0; }

  /* Upload cards */
  .upload-card {
    background: #1e293b;
    border: 1.5px dashed #334155;
    border-radius: 12px;
    padding: 20px;
    transition: border-color 0.2s;
  }
  .upload-card:hover { border-color: #60a5fa; }

  .card-label {
    font-size: 0.75rem; font-weight: 600;
    text-transform: uppercase; letter-spacing: 0.08em;
    color: #64748b; margin-bottom: 8px;
  }
  .card-title {
    font-size: 1rem; font-weight: 600; color: #e2e8f0;
    margin-bottom: 4px;
  }
  .card-desc { font-size: 0.85rem; color: #94a3b8; }

  /* File uploader overrides */
  [data-testid="stFileUploader"] section {
    border: none !important;
    background: transparent !important;
    padding: 0 !important;
  }
  [data-testid="stFileUploader"] button {
    background: #334155 !important;
    color: #e2e8f0 !important;
    border-radius: 8px !important;
    border: none !important;
    font-size: 0.85rem !important;
    padding: 6px 16px !important;
  }

  /* Merge button */
  [data-testid="stButton"] > button {
    width: 100%;
    background: linear-gradient(135deg, #3b82f6, #8b5cf6) !important;
    color: white !important;
    border: none !important;
    border-radius: 12px !important;
    font-size: 1rem !important;
    font-weight: 600 !important;
    padding: 14px !important;
    margin-top: 8px;
    transition: opacity 0.2s !important;
  }
  [data-testid="stButton"] > button:hover { opacity: 0.88 !important; }

  /* Download button */
  [data-testid="stDownloadButton"] > button {
    width: 100%;
    background: linear-gradient(135deg, #10b981, #059669) !important;
    color: white !important;
    border: none !important;
    border-radius: 12px !important;
    font-size: 1rem !important;
    font-weight: 600 !important;
    padding: 14px !important;
  }

  /* Stat cards */
  .stat-row { display: flex; gap: 16px; margin-top: 24px; }
  .stat {
    flex: 1; background: #1e293b; border-radius: 10px;
    border: 1px solid #334155;
    padding: 16px; text-align: center;
  }
  .stat-val { font-size: 1.75rem; font-weight: 700; color: #60a5fa; }
  .stat-key { font-size: 0.75rem; color: #64748b; margin-top: 2px; }

  /* Warning box */
  .warn-box {
    background: #1a1a00; border: 1px solid #ca8a04;
    border-radius: 10px; padding: 12px 16px; margin-top: 16px;
  }
  .warn-box summary { color: #fbbf24; cursor: pointer; font-size: 0.9rem; font-weight: 600; }
  .warn-box pre { color: #fde68a; font-size: 0.78rem; margin: 8px 0 0; white-space: pre-wrap; }

  /* Divider */
  hr { border-color: #1e293b; margin: 24px 0; }

  /* Step labels */
  .step {
    display: flex; align-items: center; gap: 12px;
    margin-bottom: 20px;
  }
  .step-num {
    background: #3b82f6; color: white;
    border-radius: 50%; width: 28px; height: 28px;
    display: flex; align-items: center; justify-content: center;
    font-size: 0.85rem; font-weight: 700; flex-shrink: 0;
  }
  .step-text { font-size: 0.95rem; color: #cbd5e1; }
</style>
""", unsafe_allow_html=True)

# ─── Hero ────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="hero">
  <h1>📊 Slide Merge Tool</h1>
  <p>Take the content from your <b>Root</b> file and apply the design from your <b>Template</b> file — in one click.</p>
</div>
""", unsafe_allow_html=True)

# ─── Program selector ────────────────────────────────────────────────────────
program = st.radio(
    "Choose a program",
    ["IELTS", "SEJ"],
    horizontal=True,
    help="Pick the program the files belong to — each uses its own merge rules. "
         "IELTS applies the IELTS template logic; SEJ applies the Spoken English "
         "Junior template logic.",
)

st.markdown("<hr>", unsafe_allow_html=True)

# ─── Instructions ────────────────────────────────────────────────────────────
st.markdown(f"""
<div class="step"><div class="step-num">1</div><div class="step-text">Upload your <b>{program} Root PPTX</b> — the file with all your content (text, tables, images).</div></div>
<div class="step"><div class="step-num">2</div><div class="step-text">Upload your <b>{program} Template PPTX</b> — the file with the design, theme & slide layouts you want to use.</div></div>
<div class="step"><div class="step-num">3</div><div class="step-text">Click <b>Merge Slides</b> to apply the <b>{program}</b> merge rules and download your styled presentation!</div></div>
""", unsafe_allow_html=True)

st.markdown("<hr>", unsafe_allow_html=True)

# ─── Upload columns ──────────────────────────────────────────────────────────
col1, col2 = st.columns(2)

with col1:
    st.markdown("""
    <div class="upload-card">
      <div class="card-label">Step 1</div>
      <div class="card-title">📄 Root File</div>
      <div class="card-desc">Your content source — slides with text, tables & images</div>
    </div>
    """, unsafe_allow_html=True)
    root_file = st.file_uploader(
        "root", type=["pptx"], label_visibility="collapsed", key="root"
    )
    if root_file:
        st.success(f"✅ {root_file.name}")

with col2:
    st.markdown("""
    <div class="upload-card">
      <div class="card-label">Step 2</div>
      <div class="card-title">🎨 Template File</div>
      <div class="card-desc">Your design source — theme, fonts, backgrounds & layouts</div>
    </div>
    """, unsafe_allow_html=True)
    template_file = st.file_uploader(
        "template", type=["pptx"], label_visibility="collapsed", key="template"
    )
    if template_file:
        st.success(f"✅ {template_file.name}")

st.markdown("<hr>", unsafe_allow_html=True)

# ─── Merge button ────────────────────────────────────────────────────────────
if st.button(f"✨  Merge {program} Slides", disabled=not (root_file and template_file)):
    with st.spinner(f"Merging {program} presentation... this may take a few seconds."):
        try:
            root_bytes     = root_file.read()
            template_bytes = template_file.read()
            merge_func = MERGE_FUNCS[program]
            output_bytes, warnings = merge_func(root_bytes, template_bytes)

            st.session_state["output_bytes"] = output_bytes
            st.session_state["warnings"]     = warnings
            st.session_state["root_name"]    = root_file.name

        except Exception as e:
            st.error(f"❌ Merge failed: {e}")
            st.session_state.pop("output_bytes", None)

# ─── Result ──────────────────────────────────────────────────────────────────
if "output_bytes" in st.session_state:
    out  = st.session_state["output_bytes"]
    warn = st.session_state.get("warnings", [])
    name = st.session_state.get("root_name", "presentation")
    out_name = name.replace(".pptx", "").replace("_Rootfile", "") + "_Styled.pptx"

    # Stats
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(out))
    n_slides = len(prs.slides)
    n_shapes = sum(len(s.shapes) for s in prs.slides)

    st.markdown(f"""
    <div class="stat-row">
      <div class="stat"><div class="stat-val">{n_slides}</div><div class="stat-key">Slides Merged</div></div>
      <div class="stat"><div class="stat-val">{n_shapes}</div><div class="stat-key">Total Shapes</div></div>
      <div class="stat"><div class="stat-val">{len(warn)}</div><div class="stat-key">Warnings</div></div>
    </div>
    """, unsafe_allow_html=True)

    if warn:
        items = "\n".join(warn[:20])
        if len(warn) > 20:
            items += f"\n... and {len(warn) - 20} more"
        st.markdown(f"""
        <div class="warn-box">
          <details><summary>⚠️ {len(warn)} element(s) had minor issues (non-critical)</summary>
          <pre>{items}</pre></details>
        </div>
        """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.download_button(
        label="⬇️  Download Styled Presentation",
        data=out,
        file_name=out_name,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
